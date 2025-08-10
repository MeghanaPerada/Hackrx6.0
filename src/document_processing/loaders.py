# src/document_processing/loaders.py
"""Document loading and parsing utilities"""

import os
import csv
import io
import asyncio
import requests
import pymupdf
import docx
from bs4 import BeautifulSoup
from email import policy
from email.parser import BytesParser
from pathlib import Path
from typing import List, Optional, Dict, Any
import openpyxl
from pptx import Presentation
from langchain.schema.document import Document

from ..core.config import CACHE_DIR
from ..utils.aikipedia_api import get_aikipedia_client, DEFAULT_IMAGE_DESCRIPTION_PROMPT

def download_file(url: str) -> Optional[str]:
    """Download a file from URL and return local path"""
    print(f"📥 Downloading document from {url}...")
    try:
        # Use a timeout to prevent hanging on very large files
        with requests.get(url, stream=True, timeout=300) as r: # 5 minute timeout
            r.raise_for_status()
            # Sanitize filename
            file_name = "".join(c for c in Path(url).name.split('?')[0] if c.isalnum() or c in ('.', '_', '-')).rstrip()
            if not file_name: # handle cases where URL ends in /
                import hashlib
                file_name = hashlib.md5(url.encode()).hexdigest()
            local_filename = CACHE_DIR / file_name
            with open(local_filename, 'wb') as f:
                for chunk in r.iter_content(chunk_size=8192):
                    f.write(chunk)
        print(f"✅ Download complete: {local_filename}")
        return str(local_filename)
    except requests.exceptions.RequestException as e:
        print(f"❌ Error downloading file: {e}")
        return None

def load_pdf_with_pymupdf(file_path: str, url: str) -> List[Document]:
    """Load PDF file using PyMuPDF"""
    docs = []
    try:
        pdf = pymupdf.open(file_path)
        for page_num, page in enumerate(pdf):
            text = page.get_text()
            if text.strip():
                docs.append(Document(page_content=text, metadata={"source": url, "page": page_num + 1}))
        pdf.close()
    except Exception as e:
        print(f"❌ Error loading PDF: {e}")
    return docs

def load_docx_file(file_path: str, url: str) -> List[Document]:
    """Load DOCX file"""
    try:
        document = docx.Document(file_path)
        full_text = "\n".join([para.text for para in document.paragraphs])
        if full_text.strip():
            return [Document(page_content=full_text, metadata={"source": url})]
    except Exception as e:
        print(f"❌ Error loading DOCX: {e}")
    return []

def load_xlsx_file(file_path: str, url: str) -> List[Document]:
    """Load XLSX file and convert to CSV format with advanced text cleaning"""
    from ..utils.text_cleaner import get_text_cleaner
    
    csv_sheets = []
    text_cleaner = get_text_cleaner()
    cleaning_reports = []
    
    try:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        print(f"🧹 Processing XLSX file with {len(workbook.sheetnames)} sheets for malicious content...")
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]

            output = io.StringIO()
            writer = csv.writer(output)

            has_data = False
            for row in sheet.iter_rows():
                cell_values = [cell.value if cell.value is not None else "" for cell in row]
                if any(cell_values):
                    writer.writerow(cell_values)
                    has_data = True

            if has_data:
                raw_sheet_content = f"--- Sheet: {sheet_name} ---\n" + output.getvalue()
                
                # Apply advanced text cleaning to detect and remove malicious prompts
                cleaned_content, cleaning_report = text_cleaner.clean_xlsx_content(
                    raw_sheet_content, sheet_name
                )
                
                if cleaned_content.strip():
                    csv_sheets.append(cleaned_content)
                    cleaning_reports.append(cleaning_report)
                    
                    # Log cleaning results
                    if cleaning_report.get("suspicious_rows", 0) > 0:
                        print(f"⚠️ Sheet '{sheet_name}': Found {cleaning_report['suspicious_rows']} suspicious rows")
                    else:
                        print(f"✅ Sheet '{sheet_name}': Clean content, no threats detected")

        if csv_sheets:
            full_text = "\n\n".join(csv_sheets)
            
            # Add cleaning metadata to document metadata
            metadata = {
                "source": url,
                "cleaning_applied": True,
                "sheets_processed": len(csv_sheets),
                "total_suspicious_content": sum(r.get("suspicious_rows", 0) for r in cleaning_reports),
                "cleaning_timestamp": text_cleaner.get_cleaning_stats()["last_updated"]
            }
            
            print(f"🛡️ XLSX cleaning complete: {len(csv_sheets)} sheets processed")
            print(f"📊 Cleaning stats: {text_cleaner.get_cleaning_stats()}")
            
            return [Document(page_content=full_text, metadata=metadata)]

    except Exception as e:
        print(f"❌ Error loading XLSX as CSV: {e}")
    return []

async def load_pptx_file_with_ai_analysis(file_path: str) -> List[Dict[str, Any]]:
    """Extract text and images from PPTX file with AI-powered image analysis"""
    presentation_content = []
    api_client = get_aikipedia_client()
    
    try:
        prs = Presentation(file_path)
        print(f"🎯 Processing PPTX with {len(prs.slides)} slides using AI image analysis...")
        
        for i, slide in enumerate(prs.slides):
            print(f"📄 Processing slide {i+1}/{len(prs.slides)}...")
            
            # Extract text from slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            # Process images with AI analysis
            slide_images = []
            image_descriptions = []
            
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture type
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_filename = f"slide_{i+1}_img_{len(slide_images)+1}.{image.ext}"
                        image_path = CACHE_DIR / image_filename
                        
                        # Save image locally (temporary)
                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)
                        
                        print(f"🖼️ Analyzing image: {image_filename}")
                        
                        # Process image with AI analysis and cleanup
                        description = await api_client.process_image_with_cleanup(
                            str(image_path),
                            DEFAULT_IMAGE_DESCRIPTION_PROMPT
                        )
                        
                        if description:
                            image_descriptions.append(f"Image {len(slide_images)+1}: {description}")
                            print(f"✅ Image analysis complete for {image_filename}")
                        else:
                            image_descriptions.append(f"Image {len(slide_images)+1}: [Analysis failed - image present but could not be processed]")
                            print(f"⚠️ Image analysis failed for {image_filename}")
                        
                        slide_images.append(str(image_path))
                        
                        # Clean up local image file
                        try:
                            os.unlink(image_path)
                        except:
                            pass
                            
                    except Exception as img_error:
                        print(f"❌ Error processing image in slide {i+1}: {img_error}")
                        image_descriptions.append(f"Image {len(slide_images)+1}: [Error processing image]")

            presentation_content.append({
                "slide_number": i + 1,
                "text": "\n".join(slide_text),
                "image_paths": slide_images,  # Keep for compatibility
                "image_descriptions": image_descriptions,
                "combined_content": _combine_slide_content(slide_text, image_descriptions)
            })
            
        print(f"✅ PPTX processing complete: {len(presentation_content)} slides processed")
        
    except Exception as e:
        print(f"❌ Error loading PPTX: {e}")
    
    return presentation_content

def load_pptx_file(file_path: str) -> List[Dict[str, Any]]:
    """Synchronous wrapper for PPTX loading - maintains backward compatibility"""
    try:
        # Try to get existing event loop
        loop = asyncio.get_event_loop()
        if loop.is_running():
            # If we're already in an async context, we need to handle this differently
            # For now, fall back to basic processing
            print("⚠️ Running in async context - using basic PPTX processing without AI analysis")
            return _load_pptx_file_basic(file_path)
        else:
            return loop.run_until_complete(load_pptx_file_with_ai_analysis(file_path))
    except RuntimeError:
        # No event loop exists, create one
        return asyncio.run(load_pptx_file_with_ai_analysis(file_path))

def _load_pptx_file_basic(file_path: str) -> List[Dict[str, Any]]:
    """Basic PPTX processing without AI analysis (fallback)"""
    presentation_content = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            slide_images = []
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    image = shape.image
                    image_bytes = image.blob
                    image_filename = f"slide_{i+1}_img_{len(slide_images)+1}.{image.ext}"
                    image_path = CACHE_DIR / image_filename
                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)
                    slide_images.append(str(image_path))

            presentation_content.append({
                "slide_number": i + 1,
                "text": "\n".join(slide_text),
                "image_paths": slide_images,
                "image_descriptions": [],
                "combined_content": "\n".join(slide_text)
            })
    except Exception as e:
        print(f"❌ Error loading PPTX: {e}")
    return presentation_content

def _combine_slide_content(slide_text: List[str], image_descriptions: List[str]) -> str:
    """Combine slide text and image descriptions into unified content"""
    content_parts = []
    
    # Add slide text
    if slide_text:
        text_content = "\n".join(slide_text)
        if text_content.strip():
            content_parts.append(f"Text: {text_content}")
    
    # Add image descriptions
    if image_descriptions:
        content_parts.append("Image Descriptions:")
        for desc in image_descriptions:
            content_parts.append(f"- {desc}")
    
    return "\n".join(content_parts)

def load_eml_file(file_path: str, url: str) -> List[Document]:
    """Load EML email file"""
    try:
        with open(file_path, 'rb') as f:
            msg = BytesParser(policy=policy.default).parse(f)
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    body = part.get_payload(decode=True).decode(part.get_content_charset(), 'ignore')
                    break
        elif msg.get_content_type() == "text/plain":
            body = msg.get_payload(decode=True).decode(msg.get_content_charset(), 'ignore')
        if body.strip():
            return [Document(page_content=body, metadata={"source": url, "subject": msg.get('subject', 'N/A')})]
    except Exception as e:
        print(f"❌ Error loading EML: {e}")
    return []

def load_txt_file(file_path: str, url: str) -> List[Document]:
    """Load plain text file"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        if text.strip():
            return [Document(page_content=text, metadata={"source": url})]
    except Exception as e:
        print(f"❌ Error loading TXT: {e}")
    return []

def load_csv_file(file_path: str, url: str) -> List[Document]:
    """Load CSV file"""
    docs = []
    try:
        with open(file_path, mode='r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                row_content = ", ".join(filter(None, row))
                if row_content.strip():
                    docs.append(Document(page_content=row_content, metadata={"source": url, "row": i + 1}))
    except Exception as e:
        print(f"❌ Error loading CSV: {e}")
    return docs

def load_html_file(file_path: str, url: str) -> List[Document]:
    """Load HTML file and extract text"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f, 'lxml')
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        clean_text = '\n'.join(chunk for chunk in chunks if chunk)
        if clean_text.strip():
            return [Document(page_content=clean_text, metadata={"source": url})]
    except Exception as e:
        print(f"❌ Error loading HTML: {e}")
    return []

def get_document_loader(file_ext: str):
    """Get the appropriate document loader for file extension"""
    loaders = {
        '.pdf': load_pdf_with_pymupdf,
        '.docx': load_docx_file,
        '.eml': load_eml_file,
        '.txt': load_txt_file,
        '.csv': load_csv_file,
        '.html': load_html_file,
        '.htm': load_html_file,
        '.xlsx': load_xlsx_file,
        '.pptx': load_pptx_file_wrapper
    }
    return loaders.get(file_ext)

def load_pptx_file_wrapper(file_path: str, url: str) -> List[Document]:
    """Wrapper for PPTX loading to match expected signature with AI-enhanced content"""
    try:
        presentation_content = load_pptx_file(file_path)
        documents = []
        
        for slide_data in presentation_content:
            # Use combined content (text + image descriptions) for better context
            combined_content = slide_data.get('combined_content', '')
            slide_text = slide_data.get('text', '')
            slide_number = slide_data.get('slide_number', 1)
            image_descriptions = slide_data.get('image_descriptions', [])
            
            # Prefer combined content, fall back to text only
            content = combined_content if combined_content.strip() else slide_text
            
            if content.strip():
                metadata = {
                    "source": url,
                    "slide": slide_number,
                    "images": slide_data.get('image_paths', []),
                    "has_ai_analysis": len(image_descriptions) > 0,
                    "image_count": len(slide_data.get('image_paths', []))
                }
                
                documents.append(Document(
                    page_content=content,
                    metadata=metadata
                ))
        
        print(f"📚 Created {len(documents)} document chunks from PPTX slides")
        return documents
        
    except Exception as e:
        print(f"❌ Error loading PPTX: {e}")
        return []